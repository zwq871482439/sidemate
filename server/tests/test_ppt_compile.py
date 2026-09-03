# -*- coding: utf-8 -*-
"""0.10.1 M1-E create_ppt 编译服务测试（core/ppt_compile.py）

覆盖：
- validate_svg 质量门：合法页通过 / script、class、foreignObject、错误 viewBox、
  HTML 命名实体、XML 未转义 & 各自被拒
- begin_deck：建目录 + spec_lock 落盘 + 空标题拒绝
- add_page：质量门未过不落盘（svg_output 保持永远可编译）/ 未 begin 拒绝 / 页码越界
- build_deck 全链：begin → page×2 → build → python-pptx 校验 slide 数
  （vendor ppt-master 链 + svg_quality 终检 + postflight 报告解析）
- list_decks 回放结构

注意：不 import server（看门狗）；_workspace_root monkeypatch 到 tmp_path。
"""
import json
import os

import pytest

from core import ppt_compile as pc


# PoC 验收过的 SVG 页（_ppt_poc/proj/svg_output/P01.svg，DNA-01 深蓝金）
SVG_P1 = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#0F2B46"/>
  <rect x="0" y="640" width="1280" height="80" fill="#E8B54D"/>
  <circle cx="1120" cy="180" r="140" fill="#1B4F72" opacity="0.6"/>
  <circle cx="1180" cy="240" r="60" fill="#E8B54D" opacity="0.85"/>
  <text x="90" y="300" font-family="Microsoft YaHei" font-size="72" font-weight="700" fill="#FFFFFF">桌伴 Sidemate</text>
  <text x="92" y="370" font-family="Microsoft YaHei" font-size="34" fill="#CFE8FF">本地优先的桌面 AI 工作台</text>
  <rect x="90" y="610" width="120" height="6" fill="#E8B54D"/>
</svg>
"""

SVG_P2 = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#0F2B46"/>
  <text x="90" y="140" font-family="Microsoft YaHei" font-size="54" font-weight="700" fill="#FFFFFF">核心验证点</text>
  <rect x="90" y="180" width="120" height="6" fill="#E8B54D"/>
  <text x="90" y="280" font-family="Microsoft YaHei" font-size="26" fill="#FFFFFF">中文文本可编辑性：导出后为原生文本框</text>
  <text x="90" y="340" font-family="Microsoft YaHei" font-size="26" fill="#FFFFFF">几何保真：形状位置与颜色应与预览一致</text>
  <text x="1120" y="660" font-family="Microsoft YaHei" font-size="18" fill="#9CC9EC">02 / 桌伴 PoC</text>
</svg>
"""


@pytest.fixture
def ws(tmp_path, monkeypatch):
    """把 workspace 根指到临时目录。"""
    monkeypatch.setattr(pc, "_workspace_root", lambda chat_id: str(tmp_path))
    return tmp_path


class TestValidateSvg:
    def test_valid_passes(self):
        assert pc.validate_svg(SVG_P1) == []
        assert pc.validate_svg(SVG_P2) == []

    def test_empty(self):
        assert pc.validate_svg("") == ["SVG 内容为空"]

    def test_script_rejected(self):
        svg = SVG_P1.replace("</svg>", "<script>alert(1)</script></svg>")
        assert any("script" in i for i in pc.validate_svg(svg))

    def test_class_and_style_rejected(self):
        svg = SVG_P1.replace("<rect width", '<rect class="bg" width')
        assert any("class" in i for i in pc.validate_svg(svg))
        svg2 = SVG_P1.replace("</svg>", "<style>.a{fill:red}</style></svg>")
        assert any("style" in i for i in pc.validate_svg(svg2))

    def test_foreign_object_rejected(self):
        svg = SVG_P1.replace("</svg>", "<foreignObject><p>html</p></foreignObject></svg>")
        assert any("foreignObject" in i for i in pc.validate_svg(svg))

    def test_wrong_viewbox(self):
        svg = SVG_P1.replace('viewBox="0 0 1280 720"', 'viewBox="0 0 100 100"')
        assert any("viewBox" in i for i in pc.validate_svg(svg))

    def test_named_entity_rejected(self):
        # &nbsp; 是未定义实体，XML 解析直接失败（早于实体清单检查）——
        # 两条路径都算正确拒收，模型都能拿到可操作的提示
        svg = SVG_P1.replace("桌伴 Sidemate", "桌伴&nbsp;Sidemate")
        issues = pc.validate_svg(svg)
        assert any("实体" in i or "XML 解析失败" in i for i in issues)

    def test_unescaped_ampersand(self):
        svg = SVG_P1.replace("桌伴 Sidemate", "R&D 部门")
        assert any("XML 解析失败" in i for i in pc.validate_svg(svg))

    def test_event_attr_rejected(self):
        svg = SVG_P1.replace("<circle cx=\"1120\"", "<circle onclick=\"x()\" cx=\"1120\"")
        assert any("event:onclick" in i for i in pc.validate_svg(svg))


class TestDeckFlow:
    def test_begin_creates_layout(self, ws):
        r = pc.begin_deck("t1", "产品发布演示")
        assert r["ok"] and r["deck"]
        d = os.path.join(str(ws), "ppt", r["deck"])
        assert os.path.isfile(os.path.join(d, "spec_lock.md"))
        assert os.path.isdir(os.path.join(d, "svg_output"))
        lock = open(os.path.join(d, "spec_lock.md"), encoding="utf-8").read()
        assert "viewBox: 0 0 1280 720" in lock
        assert "#E8B54D" in lock  # DNA-01 金

    def test_begin_empty_title_rejected(self, ws):
        assert pc.begin_deck("t1", "")["ok"] is False

    def test_page_requires_begin(self, ws):
        r = pc.add_page("t1", "不存在的deck", 1, SVG_P1)
        assert r["ok"] is False and r["error"] == "deck_not_found"

    def test_page_bounds(self, ws):
        deck = pc.begin_deck("t1", "t")["deck"]
        assert pc.add_page("t1", deck, 0, SVG_P1)["ok"] is False
        assert pc.add_page("t1", deck, pc.MAX_PAGES + 1, SVG_P1)["ok"] is False

    def test_rejected_page_not_written(self, ws):
        """质量门未过不落盘——svg_output 保持永远可编译。"""
        deck = pc.begin_deck("t1", "t")["deck"]
        bad = SVG_P1.replace("</svg>", "<script>x</script></svg>")
        r = pc.add_page("t1", deck, 1, bad)
        assert r["ok"] is False and r["issues"]
        svg_dir = os.path.join(str(ws), "ppt", deck, "svg_output")
        assert os.listdir(svg_dir) == []

    def test_page_ok_and_list(self, ws):
        deck = pc.begin_deck("t1", "产品发布演示")["deck"]
        assert pc.add_page("t1", deck, 1, SVG_P1)["ok"]
        assert pc.add_page("t1", deck, 2, SVG_P2)["ok"]
        decks = pc.list_decks("t1")
        assert len(decks) == 1
        assert decks[0]["pages"] == [1, 2]
        assert decks[0]["title"] == "产品发布演示"

    def test_build_empty_deck_rejected(self, ws):
        deck = pc.begin_deck("t1", "t")["deck"]
        r = pc.build_deck("t1", deck)
        assert r["ok"] is False and r["error"] == "no_pages"


class TestBuild:
    def test_full_chain(self, ws):
        """begin → page×2 → build：vendor 链 + 终检 + pptx 可被 python-pptx 打开。"""
        deck = pc.begin_deck("t1", "发布演示")["deck"]
        assert pc.add_page("t1", deck, 1, SVG_P1)["ok"]
        assert pc.add_page("t1", deck, 2, SVG_P2)["ok"]
        b = pc.build_deck("t1", deck, "测试输出")
        assert b["ok"], b.get("message")
        assert b["pages"] == 2
        pptx_path = os.path.join(str(ws), "测试输出.pptx")
        assert os.path.isfile(pptx_path)
        from pptx import Presentation
        assert len(Presentation(pptx_path).slides) == 2
        # 中文原生文本（非图片）：至少一页含文字形状
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        def _texts(shapes, acc):
            for sh in shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    _texts(sh.shapes, acc)
                elif sh.has_text_frame and sh.text_frame.text.strip():
                    acc.append(sh.text_frame.text.strip())
        acc = []
        _texts(Presentation(pptx_path).slides[0].shapes, acc)
        assert any("桌伴" in t for t in acc)
        # postflight 报告落盘
        assert os.path.isfile(os.path.join(str(ws), "ppt", deck, "validation", "测试输出.report.json"))

    def test_build_missing_deck(self, ws):
        assert pc.build_deck("t1", "没有这个deck")["ok"] is False
